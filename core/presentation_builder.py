"""
Построитель презентаций PowerPoint.

Этот модуль содержит главный оркестратор, который собирает все компоненты
вместе для генерации итоговой презентации.
"""

import logging
from pathlib import Path
from typing import Optional
from pptx import Presentation

from models import PresentationConfig, LayoutRegistry
from models.slide_types import BaseSlideConfig, YouTubeTitleSlideConfig
from io_handlers import ResourceLoader
from core import clean_markdown_for_notes
from core.placers import ImagePlacer, MediaPlacer
from config import (
    PLACEHOLDER_TITLE_IDX,
    PLACEHOLDER_SLIDE_NUM_IDX,
    PLACEHOLDER_TITLE_LAYOUT_TITLE_IDX,
    PLACEHOLDER_TITLE_LAYOUT_SLIDE_NUM_IDX,
    PLACEHOLDER_TITLE_LAYOUT_SUBTITLE_IDX,
)

logger = logging.getLogger(__name__)


class PresentationBuilder:
    """
    Оркестратор сборки презентации из конфигурации.

    Основной компонент, который:
    1. Загружает шаблон PPTX
    2. Создаёт слайды согласно конфигурации
    3. Размещает контент (текст, изображения, заметки)
    4. Сохраняет итоговый файл

    Attributes:
        layout_registry: Реестр макетов для размещения изображений.
        resource_loader: Загрузчик ресурсов (MD файлов, изображений).
        idx_title: Индекс заполнителя для заголовка.
        idx_slide_num: Индекс заполнителя для номера слайда.

    Example:
        >>> from models import LayoutRegistry
        >>> from io_handlers import ResourceLoader, PathResolver
        >>> from config import register_default_layouts
        >>>
        >>> resolver = PathResolver(Path("config.json"))
        >>> loader = ResourceLoader(resolver)
        >>> registry = LayoutRegistry()
        >>> register_default_layouts(registry)
        >>>
        >>> builder = PresentationBuilder(registry, loader)
        >>> prs = builder.build(config, Path("template.pptx"))
        >>> builder.save(prs, Path("output.pptx"))
    """

    def __init__(
        self,
        layout_registry: LayoutRegistry,
        resource_loader: ResourceLoader,
        idx_title: int = PLACEHOLDER_TITLE_IDX,
        idx_slide_num: int = PLACEHOLDER_SLIDE_NUM_IDX,
        verbose: bool = True,
    ):
        """
        Инициализация построителя.

        Args:
            layout_registry: Реестр макетов слайдов.
            resource_loader: Загрузчик ресурсов.
            idx_title: Индекс заполнителя заголовка (по умолчанию из config).
            idx_slide_num: Индекс заполнителя номера слайда.
            verbose: Выводить ли подробные сообщения о процессе.
        """
        self.layouts = layout_registry
        self.loader = resource_loader
        self.idx_title = idx_title
        self.idx_slide_num = idx_slide_num
        self.verbose = verbose

        # Инициализация классов-помощников (Composition over Inheritance)
        self.image_placer = ImagePlacer(resource_loader, layout_registry)
        self.media_placer = MediaPlacer(resource_loader)

        self._errors = []  # Список ошибок, накопленных в процессе

        logger.debug(
            f"⚙️ Инициализация PresentationBuilder: idx_title={idx_title}, idx_slide_num={idx_slide_num}"
        )
        logger.debug("🔧 ImagePlacer и MediaPlacer инициализированы")

    def build(
        self, config: PresentationConfig, template_path: Path
    ) -> Optional[Presentation]:
        """
        Главный метод сборки презентации.

        Args:
            config: Конфигурация презентации.
            template_path: Путь к файлу шаблона .pptx.

        Returns:
            Объект Presentation или None при критической ошибке.

        Raises:
            FileNotFoundError: Если шаблон не найден.
            ValueError: Если макет не найден в шаблоне.
        """
        self._errors = []  # Сброс ошибок

        logger.info(f"🚀 Начало сборки презентации из шаблона: {template_path}")
        logger.debug(f"📂 Полный путь к шаблону: {template_path.resolve()}")

        # Шаг 1: Загрузка шаблона
        logger.debug(f"� Загрузка шаблона: {template_path}")

        try:
            prs = Presentation(str(template_path))
            logger.debug(f"✅ Шаблон загружен, слайдов в мастере: {len(prs.slide_layouts)}")
        except FileNotFoundError:
            logger.error(f"❌ Шаблон не найден: {template_path}")
            raise FileNotFoundError(f"Шаблон не найден: {template_path}")
        except Exception as e:
            logger.error(f"❌ Ошибка загрузки шаблона: {e}", exc_info=True)
            raise ValueError(f"Ошибка загрузки шаблона: {e}")

        # Шаг 2: Применение workaround для PowerPoint 2013
        # (Инициализация notes_slide для всех существующих слайдов)
        for slide in prs.slides:
            _ = slide.notes_slide

        # Шаг 3: Создание слайдов
        logger.info(f"� Создание {len(config.slides)} слайдов...")
        logger.debug(f"🔍 Глобальный макет: {config.layout_name}")

        for i, slide_cfg in enumerate(config.slides, 1):
            try:
                # Определяем макет для этого слайда
                # Если в слайде указан layout_name - используем его, иначе глобальный
                current_layout_name = slide_cfg.layout_name or config.layout_name
                
                if slide_cfg.layout_name:
                    logger.debug(f"🎭 Слайд #{i}: локальный макет '{current_layout_name}' (override)")
                else:
                    logger.debug(f"🎭 Слайд #{i}: глобальный макет '{current_layout_name}'")
                
                slide_layout = self._find_layout(prs, current_layout_name)

                if not slide_layout:
                    available = [layout.name for layout in prs.slide_layouts]
                    logger.error(f"❌ Макет '{current_layout_name}' не найден. Доступные: {available}")
                    raise ValueError(
                        f"Макет '{current_layout_name}' не найден в шаблоне. "
                        f"Доступные макеты: {available}"
                    )

                self._add_slide(prs, slide_layout, slide_cfg, i)
                logger.debug(f"✅ Слайд {i} '{slide_cfg.title}' создан успешно")
            except Exception as e:
                error_msg = f"Ошибка при создании слайда {i} ('{slide_cfg.title}'): {e}"
                self._errors.append(error_msg)
                logger.error(f"❌ {error_msg}", exc_info=True)

        # Шаг 4: Сбор ошибок из placers
        image_errors = self.image_placer.get_errors()
        media_errors = self.media_placer.get_errors()
        self._errors.extend(image_errors)
        self._errors.extend(media_errors)
        
        # Шаг 5: Вывод итогов
        total_slides = len(config.slides)
        successful_slides = total_slides - len(self._errors)
        
        if self._errors:
            logger.warning(f"⚠️ Завершено с {len(self._errors)} ошибками из {total_slides} слайдов")
            for err in self._errors:
                logger.error(f"  - {err}")
        else:
            logger.info(f"✅ Презентация успешно собрана")

        logger.info(f"📊 Создано слайдов: {successful_slides}/{total_slides}")
        return prs

    def save(self, prs: Presentation, output_path: Path) -> None:
        """
        Сохраняет презентацию в файл.

        Args:
            prs: Объект презентации для сохранения.
            output_path: Путь для сохранения файла.

        Raises:
            IOError: Если не удалось сохранить файл.
        """
        try:
            logger.debug(f"🔧 Сохранение презентации: {output_path}")
            prs.save(str(output_path))
            logger.info(f"✅ Презентация сохранена: {output_path}")
        except Exception as e:
            logger.error(f"❌ Ошибка сохранения презентации: {e}", exc_info=True)
            raise IOError(f"Ошибка сохранения презентации: {e}")

    def get_errors(self) -> list:
        """
        Возвращает список ошибок, накопленных в процессе сборки.

        Returns:
            Список строк с описаниями ошибок.
        """
        return self._errors.copy()

    def _add_slide(
        self, prs: Presentation, layout, cfg: BaseSlideConfig, number: int
    ) -> None:
        """
        Добавляет один слайд в презентацию.

        Args:
            prs: Объект презентации.
            layout: Макет слайда из шаблона.
            cfg: Конфигурация слайда (BaseSlideConfig или его подклассы).
            number: Номер слайда (для отображения).
        """
        logger.info(f"📄 Обработка слайда #{number}: '{cfg.title}' (Layout: {layout.name})")
        logger.debug(f"🔍 Тип конфига: {type(cfg).__name__}, изображений: {len(cfg.images) if cfg.images else 0}, аудио: {bool(cfg.audio)}")
        
        # Создание слайда
        slide = prs.slides.add_slide(layout)
        logger.debug(f"🔧 Слайд создан, ID: {slide.slide_id}")

        # Workaround для PowerPoint 2013
        _ = slide.notes_slide

        # Определяем, используется ли TitleLayout
        is_title_layout = isinstance(cfg, YouTubeTitleSlideConfig)

        # Выбираем правильные индексы в зависимости от типа макета
        if is_title_layout:
            idx_title = PLACEHOLDER_TITLE_LAYOUT_TITLE_IDX
            idx_slide_num = PLACEHOLDER_TITLE_LAYOUT_SLIDE_NUM_IDX
        else:
            idx_title = self.idx_title
            idx_slide_num = self.idx_slide_num

        # 1. Заголовок
        try:
            title_ph = slide.shapes.placeholders[idx_title]
            title_ph.text_frame.text = cfg.title
            logger.debug(f"🔧 Title установлен в placeholder idx={idx_title}")
        except KeyError:
            # ИЗМЕНЕНИЕ: Не падаем, если это слайд без заголовка (например, Shorts)
            logger.debug(f"⚠️ Заполнитель заголовка idx={idx_title} не найден (пропуск для графического слайда)")

        # 2. Дополнительные поля для YouTubeTitleSlideConfig
        if is_title_layout:
            logger.debug("🔧 Обработка YouTube-титульника")
            self._set_youtube_title_fields(slide, cfg)

        # 3. Номер слайда
        try:
            num_ph = slide.shapes.placeholders[idx_slide_num]
            num_ph.text_frame.text = str(number)
            logger.debug(f"🔧 Номер слайда {number} установлен в placeholder idx={idx_slide_num}")
        except KeyError:
            logger.debug(f"🔍 Заполнитель номера idx={idx_slide_num} не найден (не критично)")

        # 4. Заметки докладчика
        logger.debug(f"📝 Загрузка заметок: {cfg.notes_source}")
        notes_text = self.loader.load_notes(cfg.notes_source)
        clean_notes = clean_markdown_for_notes(notes_text)
        slide.notes_slide.notes_text_frame.text = clean_notes
        logger.debug(f"🔧 Заметки добавлены: {len(clean_notes)} символов")

        # 5. Изображения (делегируем ImagePlacer)
        logger.debug(f"🖼️ Размещение изображений: {len(cfg.images) if cfg.images else 0}")
        self.image_placer.place_images(slide, cfg)

        # 6. Аудио (делегируем MediaPlacer)
        if cfg.audio:
            logger.debug(f"🎵 Добавление аудио: {cfg.audio}")
            self.media_placer.place_audio(slide, cfg.audio)

    def _set_youtube_title_fields(self, slide, cfg: YouTubeTitleSlideConfig) -> None:
        """
        Устанавливает специфичные поля для YouTubeTitleSlideConfig.

        Args:
            slide: Объект слайда.
            cfg: Конфигурация титульного слайда YouTube.

        Note:
            Индексы заполнителей для TitleLayout в youtube_base.pptx:
            - idx=10: title (основной заголовок)
            - idx=12: slide_number (номер слайда)
            - idx=13: subtitle (подзаголовок/описание серии)
        """
        logger.debug(f"🔧 YouTube поля: subtitle='{cfg.subtitle}', series_number={cfg.series_number}")
        
        # Subtitle (placeholder idx=13 в TitleLayout)
        try:
            subtitle_ph = slide.shapes.placeholders[
                PLACEHOLDER_TITLE_LAYOUT_SUBTITLE_IDX
            ]
            subtitle_ph.text_frame.text = cfg.subtitle
            logger.debug(f"🔧 Subtitle установлен в placeholder idx={PLACEHOLDER_TITLE_LAYOUT_SUBTITLE_IDX}")
        except KeyError as e:
            logger.warning(f"⚠️ Заполнитель subtitle idx={PLACEHOLDER_TITLE_LAYOUT_SUBTITLE_IDX} не найден: {e}")
        except Exception as e:
            logger.error(f"❌ Ошибка при заполнении subtitle: {e}", exc_info=True)

        # Series number - пока нет заполнителя в шаблоне
        if cfg.series_number:
            logger.debug(f"🔍 Series number '{cfg.series_number}' не добавлен (нет заполнителя)")

    @staticmethod
    def _find_layout(prs: Presentation, layout_name: str):
        """
        Ищет макет по имени в шаблоне.

        Args:
            prs: Объект презентации.
            layout_name: Имя макета для поиска.

        Returns:
            Объект макета или None, если не найден.
        """
        logger.debug(f"🔍 Поиск макета '{layout_name}' в мастере...")
        
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                logger.debug(f"🔧 Макет '{layout_name}' найден")
                return layout
        
        logger.warning(f"⚠️ Макет '{layout_name}' не найден в шаблоне")
        return None
