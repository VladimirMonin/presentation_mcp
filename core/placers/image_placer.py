"""
Размещение изображений на слайдах.

Модуль отвечает за вставку и масштабирование изображений согласно
чертежам макетов (Layout Blueprints).
"""

import logging
from pathlib import Path
from pptx.util import Cm

from models import LayoutRegistry
from models.slide_types import BaseSlideConfig, YouTubeTitleSlideConfig
from io_handlers import ResourceLoader
from core import calculate_smart_dimensions, convert_webp_to_png

logger = logging.getLogger(__name__)


class ImagePlacer:
    """
    Класс для размещения изображений на слайдах.

    Этот класс инкапсулирует всю логику работы с изображениями:
    - Резолвинг путей через ResourceLoader
    - Автоматическая конвертация WebP → PNG (in-memory)
    - Умное масштабирование под макет (calculate_smart_dimensions)
    - Размещение согласно Blueprint из LayoutRegistry

    Attributes:
        loader: ResourceLoader для разрешения путей к изображениям.
        layouts: LayoutRegistry для получения чертежей макетов.
        errors: Список ошибок, возникших при размещении изображений.

    Example:
        >>> from models import LayoutRegistry
        >>> from io_handlers import ResourceLoader, PathResolver
        >>> from config import register_default_layouts
        >>>
        >>> resolver = PathResolver(Path("config.json"))
        >>> loader = ResourceLoader(resolver)
        >>> registry = LayoutRegistry()
        >>> register_default_layouts(registry)
        >>>
        >>> placer = ImagePlacer(loader, registry)
        >>> placer.place_images(slide, slide_config)
    """

    def __init__(self, resource_loader: ResourceLoader, layout_registry: LayoutRegistry):
        """
        Инициализация ImagePlacer.

        Args:
            resource_loader: Объект для разрешения путей к ресурсам.
            layout_registry: Реестр макетов для получения чертежей.
        """
        self.loader = resource_loader
        self.layouts = layout_registry
        self.errors = []
        logger.debug("⚙️ ImagePlacer инициализирован")

    def place_images(self, slide, cfg: BaseSlideConfig) -> bool:
        """
        Размещает все изображения для слайда согласно макету.

        Args:
            slide: Объект слайда python-pptx.
            cfg: Конфигурация слайда (BaseSlideConfig или его подклассы).

        Returns:
            True если хотя бы одно изображение успешно размещено, False если нет изображений.

        Raises:
            KeyError: Если указанный layout_type не зарегистрирован.
        """
        if not cfg.images:
            logger.debug("🔍 Нет изображений для размещения")
            return False

        logger.info(f"🖼️ Размещение изображений для слайда: '{cfg.title}'")

        # Получаем чертёж макета
        # Для YouTubeTitleSlideConfig используем фиксированный макет title_youtube
        if isinstance(cfg, YouTubeTitleSlideConfig):
            layout_type = "title_youtube"
            logger.debug("🔍 YouTube титульник -> макет 'title_youtube'")
        else:
            layout_type = cfg.layout_type
            logger.debug(f"🔍 Используем макет из конфига: '{layout_type}'")

        try:
            blueprint = self.layouts.get(layout_type)
            logger.debug(
                f"🔍 Чертеж макета '{layout_type}': требуется {blueprint.required_images} изображений"
            )
        except KeyError:
            logger.error(
                f"❌ Макет '{layout_type}' не зарегистрирован. Доступные: {self.layouts.list_all()}"
            )
            raise KeyError(
                f"Макет '{layout_type}' не зарегистрирован. "
                f"Доступные: {self.layouts.list_all()}"
            )

        # Проверка количества изображений
        if len(cfg.images) < blueprint.required_images:
            logger.warning(
                f"⚠️ Ожидалось {blueprint.required_images} изображений, предоставлено {len(cfg.images)}"
            )

        # Размещение каждого изображения
        success_count = 0
        for i, img_path_str in enumerate(cfg.images):
            if i >= len(blueprint.placements):
                # Больше изображений, чем размещений - игнорируем лишние
                logger.warning(
                    f"⚠️ Изображение #{i + 1} '{img_path_str}' игнорируется (нет размещения в макете)"
                )
                break

            if self._place_single_image(slide, img_path_str, blueprint.placements[i]):
                success_count += 1

        logger.debug(f"📊 Размещено изображений: {success_count}/{len(cfg.images)}")
        return success_count > 0

    def _place_single_image(self, slide, img_path_str: str, placement) -> bool:
        """
        Размещает одно изображение на слайде.

        Args:
            slide: Объект слайда python-pptx.
            img_path_str: Путь к изображению (строка).
            placement: Объект ImagePlacement с координатами и размерами.

        Returns:
            True если изображение успешно размещено, False в случае ошибки.
        """
        try:
            logger.debug(f"📍 Размещение изображения: {img_path_str}")

            # 1. Разрешение пути к изображению
            img_path = self.loader.resolve_image(img_path_str)

            # 2. Автоматическая конвертация WebP → PNG (in-memory)
            original_path = img_path
            image_source = img_path  # По умолчанию используем путь к файлу

            if img_path.suffix.lower() == ".webp":
                try:
                    # convert_webp_to_png теперь возвращает BytesIO
                    image_source = convert_webp_to_png(img_path)
                    logger.debug(f"🔄 WebP сконвертирован в памяти: {original_path.name}")
                except Exception as e:
                    error_msg = f"Ошибка конвертации WebP {img_path_str}: {e}"
                    self.errors.append(error_msg)
                    logger.error(f"❌ {error_msg}", exc_info=True)
                    return False

            # 3. Получение параметров размещения
            placement_dict = placement.to_dict()

            logger.debug(
                f"📏 Чертеж: left={placement_dict['left']}, top={placement_dict['top']}, "
                f"max_width={placement_dict['max_width']}, max_height={placement_dict['max_height']}"
            )

            # 4. Умное масштабирование (для BytesIO используем исходный путь)
            dimensions_source = (
                original_path if img_path.suffix.lower() == ".webp" else img_path
            )
            width, height = calculate_smart_dimensions(
                dimensions_source,
                placement_dict["max_width"],
                placement_dict["max_height"],
            )

            # 5. Конвертация в единицы python-pptx
            left_cm = Cm(placement_dict["left"])
            top_cm = Cm(placement_dict["top"])
            width_cm = Cm(width) if width is not None else None
            height_cm = Cm(height) if height is not None else None

            width_str = f"{width:.2f}" if width is not None else "auto"
            height_str = f"{height:.2f}" if height is not None else "auto"
            logger.debug(
                f"📐 Вычислено (см): left={placement_dict['left']:.2f}, top={placement_dict['top']:.2f}, "
                f"w={width_str}, h={height_str}"
            )

            # EMU для детального логирования
            emu_left = int(left_cm)
            emu_top = int(top_cm)
            emu_width = int(width_cm) if width_cm else None
            emu_height = int(height_cm) if height_cm else None

            logger.debug(
                f"🎯 Финальные EMU: left={emu_left}, top={emu_top}, "
                f"width={emu_width or 'auto'}, height={emu_height or 'auto'}"
            )

            # 6. Добавление изображения на слайд
            # python-pptx поддерживает как пути (str/Path), так и потоки (BytesIO)
            if isinstance(image_source, Path):
                slide.shapes.add_picture(
                    str(image_source),
                    left_cm,
                    top_cm,
                    width=width_cm,
                    height=height_cm,
                )
            else:
                # BytesIO передаём напрямую
                slide.shapes.add_picture(
                    image_source, left_cm, top_cm, width=width_cm, height=height_cm
                )

            logger.debug(f"✅ Изображение '{img_path_str}' успешно размещено")
            return True

        except FileNotFoundError:
            # Изображение не найдено - добавляем в ошибки, но продолжаем
            error_msg = f"Изображение не найдено: {img_path_str}"
            self.errors.append(error_msg)
            logger.warning(f"⚠️ {error_msg}")
            return False

        except Exception as e:
            # Другая ошибка при добавлении изображения
            error_msg = f"Ошибка добавления изображения {img_path_str}: {e}"
            self.errors.append(error_msg)
            logger.error(f"❌ {error_msg}", exc_info=True)
            return False

    def get_errors(self) -> list:
        """
        Возвращает список ошибок, накопленных в процессе работы.

        Returns:
            Копия списка строк с описаниями ошибок.
        """
        return self.errors.copy()

    def clear_errors(self) -> None:
        """Очищает список накопленных ошибок."""
        self.errors = []
        logger.debug("🧹 Список ошибок ImagePlacer очищен")
