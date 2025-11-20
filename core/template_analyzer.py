"""
Анализатор шаблонов PowerPoint.

Этот модуль позволяет анализировать структуру .pptx шаблонов
для определения индексов заполнителей и параметров макетов.
"""

import logging
from pathlib import Path
from pptx import Presentation

logger = logging.getLogger(__name__)


def analyze_template(template_path: Path, layout_name: str = "VideoLayout") -> None:
    """
    Анализирует шаблон PPTX и выводит информацию о макетах и заполнителях.

    Args:
        template_path: Путь к файлу шаблона .pptx.
        layout_name: Имя макета для детального анализа (опционально).

    Example:
        >>> analyze_template(Path("template.pptx"), "VideoLayout")
    """
    logger.info(f"🔍 Анализ шаблона: {template_path}")
    logger.debug(f"🔍 Целевой макет для детального анализа: '{layout_name}'")
    
    try:
        prs = Presentation(str(template_path))
        logger.debug("🔧 Шаблон успешно загружен")
    except FileNotFoundError:
        logger.error(f"❌ Файл не найден: {template_path}")
        return
    except Exception as e:
        logger.error(f"❌ Ошибка загрузки шаблона: {e}", exc_info=True)
        return

    # Вывод всех доступных макетов
    logger.info(f"📋 Найдено макетов: {len(prs.slide_layouts)}")
    
    for i, layout in enumerate(prs.slide_layouts):
        logger.debug(f"🔍 Макет #{i + 1}: '{layout.name}'")

    # Поиск нужного макета
    target_layout = None
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            target_layout = layout
            break

    if not target_layout:
        logger.warning(f"⚠️ Макет '{layout_name}' не найден в шаблоне")
        return

    # Детальный анализ макета
    logger.debug(f"🔍 Начинаем детальный анализ макета '{layout_name}'")

    placeholders = target_layout.placeholders

    if not placeholders:
        logger.debug(f"🔍 Макет '{layout_name}' не содержит заполнителей")
        return

    logger.debug(f"📋 Найдено заполнителей: {len(placeholders)}")

    for ph in placeholders:
        logger.debug(
            f"🔧 Заполнитель: idx={ph.placeholder_format.idx}, "
            f"type={ph.placeholder_format.type}, name='{ph.name}'"
        )

        # Попытка получить текст (если есть)
        try:
            if hasattr(ph, "text_frame") and ph.text_frame:
                sample_text = (
                    ph.text_frame.text[:50] if ph.text_frame.text else "(пусто)"
                )
                logger.debug(f"🔧 Текст заполнителя: {sample_text}")
        except Exception:
            pass

    logger.info("✅ Анализ шаблона завершён успешно")


def list_layouts(template_path: Path) -> None:
    """
    Выводит простой список макетов в шаблоне.

    Args:
        template_path: Путь к файлу шаблона.
    """
    logger.info(f"📋 Запрошен список макетов для: {template_path}")
    
    try:
        prs = Presentation(str(template_path))
        logger.debug(f"🔧 Шаблон успешно загружен: {len(prs.slide_layouts)} макетов")
    except Exception as e:
        logger.error(f"❌ Ошибка загрузки шаблона: {e}", exc_info=True)
        return

    for i, layout in enumerate(prs.slide_layouts, 1):
        logger.debug(f"🔍 Макет #{i}: '{layout.name}'")
    
    logger.info(f"✅ Список макетов выведен: {len(prs.slide_layouts)} шт.")
