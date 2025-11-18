# pip install python-pptx
from pptx import Presentation
import sys

# --- НАСТРОЙКИ ---
# Имя твоего шаблона
TEMPLATE_FILE = 'template.pptx' 
# Имя твоего главного макета (которое ты задал в Образце слайдов)
TARGET_LAYOUT_NAME = 'VideoLayout' 
# --------------------

def analyze_layout(template_path, layout_name):
    """
    Загружает шаблон и выводит все заполнители (placeholders)
    и их 'idx' (индексы) для конкретного макета (layout).
    """
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Ошибка загрузки шаблона '{template_path}': {e}")
        sys.exit(1)

    # 1. Находим наш макет по имени (НАДЕЖНЫЙ СПОСОБ)
    target_layout = None
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            target_layout = layout
            break

    if not target_layout:
        print(f"Ошибка: Макет с именем '{layout_name}' не найден в шаблоне.")
        print("Доступные макеты:")
        for layout in prs.slide_layouts:
            print(f"- {layout.name}")
        sys.exit(1)

    print(f"--- Анализ макета: '{target_layout.name}' ---")

    # 2. Перебираем все заполнители на этом макете
    try:
        placeholders = target_layout.placeholders
    except Exception as e:
        print(f"Не удалось получить заполнители: {e}")
        return

    if not placeholders:
        print("На этом макете не найдено ни одного заполнителя (placeholder).")
        return

    print(f"Найдено {len(placeholders)} заполнителей:\n")
    
    # 3. Выводим их 'idx' и 'name'
    for ph in placeholders:
        # 'idx' хранится внутри 'placeholder_format'
        idx = ph.placeholder_format.idx
        name = ph.name
        
        print(f"Имя: {name}")
        print(f"  idx: {idx}") # Это та цифра, которая тебе нужна
        print("  ---")

    print("\nИспользуй эти 'idx' в своем основном скрипте для доступа")
    print("через: slide.shapes.placeholders[idx_number]")

# --- Запуск анализатора ---
if __name__ == "__main__":
    analyze_layout(TEMPLATE_FILE, TARGET_LAYOUT_NAME)