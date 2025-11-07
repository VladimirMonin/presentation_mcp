from pptx import Presentation
from pptx.util import Inches, Cm  # <<<--- ИМПОРТИРУЕМ Cm (Сантиметры)
import json
import sys
from two import clean_markdown_for_notes  # Импортируем наш очиститель

try:
    from PIL import Image  # Импортируем Pillow
except ImportError:
    print("Ошибка: Библиотека Pillow не найдена.")
    print("Пожалуйста, установите ее: pip install Pillow")
    sys.exit(1)


# -----------------------------------------------------------------
# --- "КОНТРАКТ API" (Получено из analyze_template.py) ---
# -----------------------------------------------------------------
TEMPLATE_FILE = "template.pptx"
TARGET_LAYOUT_NAME = "VideoLayout"  # Имя макета из "Образца слайдов"

# ID Заполнителей (ТВОИ РЕАЛЬНЫЕ ID)
IDX_TITLE = 10
IDX_SLIDE_NUM = 11

# -----------------------------------------------------------------
# --- "ЧЕРТЕЖИ" МАКЕТОВ (layout_type) ---
# --- ТВОИ ТОЧНЫЕ СПЕЦИФИКАЦИИ В Cm ---
# -----------------------------------------------------------------

# --- Макет 'single_wide' (Слайд 1: "Качаем VS Code") ---
WIDE_LEFT = Cm(10.2)
WIDE_TOP = Cm(4.2)
WIDE_MAX_WIDTH = Cm(20)
WIDE_MAX_HEIGHT = Cm(10)

# --- Макет 'single_tall' (Слайд 2: "Обзор расширений") ---
TALL_LEFT = Cm(10.46)
TALL_TOP = Cm(2.96)
TALL_MAX_WIDTH = Cm(11.2)
TALL_MAX_HEIGHT = Cm(15.2)

# --- Макет 'two_stack' (Слайд 3: "Cline и Excalidraw") ---
# 1-я (верхняя) картинка
STACK_1_LEFT = Cm(10.16)
STACK_1_TOP = Cm(3.47)
STACK_1_MAX_WIDTH = Cm(18.4)
STACK_1_MAX_HEIGHT = Cm(3.91)

# 2-я (нижняя) картинка
STACK_2_LEFT = Cm(10.16)
STACK_2_TOP = Cm(11)
STACK_2_MAX_WIDTH = Cm(18.07)
STACK_2_MAX_HEIGHT = Cm(4.58)

# -----------------------------------------------------------------


def load_slide_data(json_path: str) -> list:
    """Загружает данные о слайдах из JSON-файла."""
    try:
        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return data.get("slides", [])
    except FileNotFoundError:
        print(f"Ошибка: JSON-файл не найден: {json_path}")
        return []
    except json.JSONDecodeError:
        print(f"Ошибка: Некорректный JSON в файле: {json_path}")
        return []


def get_smart_dimensions(
    image_path, max_width_cm, max_height_cm
):  # <<-- Теперь принимаем Cm
    """
    Вычисляет, ограничить ширину или высоту,
    чтобы сохранить пропорции и влезть в Bounding Box.
    """
    try:
        with Image.open(image_path) as img:
            img_width, img_height = img.size
    except FileNotFoundError:
        print(f"  Ошибка: get_smart_dimensions не нашел файл: {image_path}")
        return None, None  # Возвращаем None, None
    except Exception as e:
        print(f"  Ошибка Pillow при чтении {image_path}: {e}")
        return None, None

    # Конвертируем "коробку" в Cm (1 дюйм = 2.54 см)
    # Для расчета пропорций это не так важно, но для ясности
    box_width = max_width_cm
    box_height = max_height_cm

    # Пропорции картинки и "коробки"
    img_ratio = img_width / img_height
    box_ratio = box_width / box_height

    if img_ratio > box_ratio:
        # Картинка ШИРЕ, чем коробка. Ограничиваем по ШИРИНЕ.
        width = max_width_cm  # <<-- Теперь это Cm
        height = None  # Пусть pptx вычислит высоту
    else:
        # Картинка ВЫШЕ, чем коробка (или такая же). Ограничиваем по ВЫСОТЕ.
        width = None  # Пусть pptx вычислит ширину
        height = max_height_cm  # <<-- Теперь это Cm

    return width, height


def generate_presentation(slide_data: list, template_path: str, output_path: str):
    """
    Основная функция генерации презентации.
    """
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Критическая ошибка: Не могу загрузить шаблон '{template_path}'. {e}")
        return

    # Ищем наш кастомный макет по имени
    slide_layout = next(
        (layout for layout in prs.slide_layouts if layout.name == TARGET_LAYOUT_NAME),
        None,
    )
    if not slide_layout:
        print(f"Критическая ошибка: Макет '{TARGET_LAYOUT_NAME}' не найден в шаблоне.")
        return

    # --- Применение "Обходного пути" бага PowerPoint 2013 ---
    for slide in prs.slides:
        _ = slide.notes_slide

    # --- Главный цикл создания слайдов ---
    for i, item in enumerate(slide_data):
        print(f"Создание слайда {i + 1}...")
        slide = prs.slides.add_slide(slide_layout)

        # --- Обходной путь бага 2013 (для *каждого* нового слайда) ---
        _ = slide.notes_slide

        # --- 1. Заполнение текста (БЕЗ 'try...except' ДЛЯ ОТЛАДКИ) ---
        # Если здесь ошибка, скрипт УПАДЕТ и покажет ее

        # Заголовок
        print(f"  ...Ищу placeholder с idx: {IDX_TITLE} (для Title)")
        title_ph = slide.shapes.placeholders[IDX_TITLE]
        title_ph.text_frame.text = item.get("title", "Без заголовка")
        print(f"  ...Заголовок '{item.get('title')}' вставлен.")

        # Номер слайда (просто как текст)
        print(f"  ...Ищу placeholder с idx: {IDX_SLIDE_NUM} (для Номера)")
        num_ph = slide.shapes.placeholders[IDX_SLIDE_NUM]
        num_ph.text_frame.text = str(i + 1)  # Авто-нумерация
        print(f"  ...Номер '{i + 1}' вставлен.")

        # --- 2. Добавление заметок докладчика (с очисткой) ---
        notes_text = item.get("notes_text", "")
        print("  ...Очищаю заметки...")
        clean_notes = clean_markdown_for_notes(notes_text)
        slide.notes_slide.notes_text_frame.text = clean_notes
        print("  ...Заметки вставлены.")

        # --- 3. Размещение скриншотов (БЕЗ 'try...except' ДЛЯ ОТЛАДКИ) ---
        layout_type = item.get("layout_type", "single_wide")  # По умолчанию single_wide
        images = item.get("images", [])

        if layout_type == "single_wide" and images:
            img_path = images[0]
            print(f"  ...Вычисляю 'smart' размеры для (single_wide): {img_path}")

            width, height = get_smart_dimensions(
                img_path, WIDE_MAX_WIDTH, WIDE_MAX_HEIGHT
            )

            if width is not None or height is not None:
                print(
                    f"  ...Добавляю картинку (single_wide): {img_path} (w:{width}, h:{height})"
                )
                slide.shapes.add_picture(
                    img_path, WIDE_LEFT, WIDE_TOP, width=width, height=height
                )
                print("  ...Картинка 'single_wide' добавлена.")
            else:
                print(f"  ОШИБКА: Не смог добавить картинку {img_path}, пропуск.")

        elif layout_type == "single_tall" and images:
            img_path = images[0]
            print(f"  ...Вычисляю 'smart' размеры для (single_tall): {img_path}")

            width, height = get_smart_dimensions(
                img_path, TALL_MAX_WIDTH, TALL_MAX_HEIGHT
            )

            if width is not None or height is not None:
                print(
                    f"  ...Добавляю картинку (single_tall): {img_path} (w:{width}, h:{height})"
                )
                slide.shapes.add_picture(
                    img_path, TALL_LEFT, TALL_TOP, width=width, height=height
                )
                print("  ...Картинка 'single_tall' добавлена.")
            else:
                print(f"  ОШИБКА: Не смог добавить картинку {img_path}, пропуск.")

        elif layout_type == "two_stack" and len(images) >= 2:
            img_path_1 = images[0]
            img_path_2 = images[1]

            # 1-я картинка
            print(f"  ...Вычисляю 'smart' размеры для (stack 1): {img_path_1}")
            width1, height1 = get_smart_dimensions(
                img_path_1, STACK_1_MAX_WIDTH, STACK_1_MAX_HEIGHT
            )
            if width1 is not None or height1 is not None:
                print(
                    f"  ...Добавляю картинку 1 (stack): {img_path_1} (w:{width1}, h:{height1})"
                )
                slide.shapes.add_picture(
                    img_path_1, STACK_1_LEFT, STACK_1_TOP, width=width1, height=height1
                )
                print("  ...Картинка 1 'stack' добавлена.")
            else:
                print(f"  ОШИБКА: Не смог добавить картинку {img_path_1}, пропуск.")

            # 2-я картинка
            print(f"  ...Вычисляю 'smart' размеры для (stack 2): {img_path_2}")
            width2, height2 = get_smart_dimensions(
                img_path_2, STACK_2_MAX_WIDTH, STACK_2_MAX_HEIGHT
            )
            if width2 is not None or height2 is not None:
                print(
                    f"  ...Добавляю картинку 2 (stack): {img_path_2} (w:{width2}, h:{height2})"
                )
                slide.shapes.add_picture(
                    img_path_2, STACK_2_LEFT, STACK_2_TOP, width=width2, height=height2
                )
                print("  ...Картинка 2 'stack' добавлена.")
            else:
                print(f"  ОШИБКА: Не смог добавить картинку {img_path_2}, пропуск.")

        else:
            if images:
                print(
                    f"Предупреждение: Неизвестный layout_type '{layout_type}' или неверное кол-во картинок."
                )

    # --- 4. Сохранение результата ---
    try:
        prs.save(output_path)
        print(f"\nПрезентация успешно сохранена: '{output_path}'")
    except Exception as e:
        print(f"Ошибка сохранения файла: {e}")


# --- Запуск ---
if __name__ == "__main__":
    # Входной JSON-конфиг (от ИИ-агента или созданный вручную)
    INPUT_JSON = "slides_config.json"
    OUTPUT_FILE = "generated_presentation2.pptx"  # Твое имя файла

    slide_data = load_slide_data(INPUT_JSON)
    if slide_data:
        generate_presentation(slide_data, TEMPLATE_FILE, OUTPUT_FILE)
